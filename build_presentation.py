"""Build the IMAGE 2026 oral presentation from the official .potx template.

Design rules, from the IMAGE26 speaker manual and the dataviz procedure:
  - 16:9, 1920x1080, PowerPoint 2016 .pptx
  - <= 10-15 words on a slide; explanation lives in the speaker notes
  - figures, not tables; every result slide is a chart
  - emphasis palette (validated): accent #2a78d6, warm #eb6834, rest gray

Charts come from figures/slide_charts.py. Regenerate those first.
"""
import os
import re
import zipfile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = '/home/carlos/transfer/IMAGE26_PPT_template.potx'
WORK = os.path.join(ROOT, '.template_base.pptx')
OUT = os.path.join(ROOT, 'IMAGE2026_FL_Seismic.pptx')
SLIDES = os.path.join(ROOT, 'figures', 'slides')

TITLE, TITLE_CONTENT, SECTION, TITLE_ONLY, BLANK = 0, 1, 2, 5, 6

# Institutional logos. Drop the real files at these paths and rebuild — they
# are placed automatically: large on the title slide, small bottom-right
# elsewhere. Missing files are skipped silently, so the deck still builds.
# Prefer PNG with transparency (SVG is not supported by python-pptx).
LOGO_DIR = os.path.join(ROOT, 'figures', 'logos')
LOGO_OLIVES = os.path.join(LOGO_DIR, 'olives.png')
LOGO_GATECH = os.path.join(LOGO_DIR, 'gatech.png')
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
WARM = RGBColor(0xEB, 0x68, 0x34)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)


def potx_to_pptx(src, dst):
    """python-pptx refuses .potx; rewrite the content type.

    The template also ships a stub slide1. Dropping it from the package (not
    merely unlinking it) keeps one part per slide — a duplicated slide1.xml
    makes PowerPoint refuse the file.
    """
    drop = {'ppt/slides/slide1.xml', 'ppt/slides/_rels/slide1.xml.rels'}
    zin = zipfile.ZipFile(src)
    with zipfile.ZipFile(dst, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            if item.filename in drop:
                continue
            data = zin.read(item.filename)
            if item.filename == '[Content_Types].xml':
                data = data.replace(b'presentationml.template.main+xml',
                                    b'presentationml.presentation.main+xml')
                data = re.sub(rb'<Override PartName="/ppt/slides/slide1\.xml"'
                              rb'[^>]*/>', b'', data)
            elif item.filename == 'ppt/presentation.xml':
                data = re.sub(rb'<p:sldIdLst>.*?</p:sldIdLst>', b'', data,
                              flags=re.S)
            elif item.filename == 'ppt/_rels/presentation.xml.rels':
                data = re.sub(rb'<Relationship[^>]*Target="slides/slide1\.xml"'
                              rb'[^>]*/>', b'', data)
            zout.writestr(item, data)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _logo(slide, path, height, right, bottom):
    """Place one logo with its bottom-right corner at (right, bottom).

    Aspect ratio is preserved from the file; only height is specified, so a
    wide wordmark and a square mark both sit on the same baseline.
    """
    if not os.path.exists(path):
        return 0
    iw, ih = Image.open(path).size
    h = height
    w = int(iw * (h / ih))
    slide.shapes.add_picture(path, right - w, bottom - h, w, h)
    return w


def add_logos(prs, slide, large=False):
    """OLIVES + Georgia Tech, bottom-right. Larger on the title slide."""
    h = Inches(0.85 if large else 0.42)
    gap = Inches(0.28 if large else 0.18)
    right = prs.slide_width - Inches(0.55)
    bottom = prs.slide_height - Inches(0.45)
    w = _logo(slide, LOGO_GATECH, h, right, bottom)
    if w:
        right -= w + gap
    _logo(slide, LOGO_OLIVES, h, right, bottom)


def _fit(prs, slide, image, top, max_h):
    """Centre an image, scaled to fit within the slide width and max_h."""
    iw, ih = Image.open(image).size
    max_w = prs.slide_width - Inches(1.6)
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(image, int((prs.slide_width - w) / 2), top, w, h)
    return h


def chart_slide(prs, title, chart, note, takeaway=None):
    """Title + chart + one-line takeaway. The workhorse layout."""
    s = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    s.shapes.title.text = title
    # Title placeholder runs to 1.85in; start the chart clear of it.
    top = Inches(2.0)
    h = _fit(prs, s, os.path.join(SLIDES, chart), top,
             Inches(4.25 if takeaway else 5.0))
    if takeaway:
        # Stop clear of the logo block in the bottom-right corner.
        box = s.shapes.add_textbox(Inches(0.8), top + h + Inches(0.18),
                                   prs.slide_width - Inches(4.4), Inches(0.7))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = takeaway
        p.alignment = 1
        r = p.runs[0]
        r.font.size = Pt(21)
        r.font.color.rgb = INK2
    notes(s, note)
    return s


def statement(prs, line1, line2, note, accent=ACCENT):
    """A full-bleed statement slide. Used to mark the turns in the argument."""
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    box = s.shapes.add_textbox(Inches(1.0), Inches(2.5),
                               prs.slide_width - Inches(2.0), Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = line1
    p.runs[0].font.size = Pt(54)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = INK
    if line2:
        p2 = tf.add_paragraph()
        p2.text = line2
        p2.runs[0].font.size = Pt(32)
        p2.runs[0].font.color.rgb = accent
    notes(s, note)
    return s


def build():
    potx_to_pptx(TEMPLATE, WORK)
    prs = Presentation(WORK)
    assert len(prs.slides) == 0, 'template stub slide should already be gone'

    # 1 ─ Title
    s = prs.slides.add_slide(prs.slide_layouts[TITLE])
    s.shapes.title.text = 'The Effectiveness of Federated Learning in Seismic Interpretation'
    s.placeholders[1].text = 'Carlos Marí · Ghassan AlRegib · OLIVES, Georgia Tech'
    notes(s, "Memorise this opening; look at the audience, not the screen.\n\n"
             "'Federated learning promises a shared seismic interpreter that "
             "no one has to hand over their data for. Our paper asked whether "
             "it actually works. The answer was: not under real geography. "
             "Today I want to show you what we have learned since — including "
             "where we were wrong.'\n\n"
             "TIMING: 15-17 min talk, green light until 5 min remain.")

    # 2 ─ The promise
    statement(prs, 'Train together.', 'Share no data.',
              "Sixty seconds of motivation. Seismic data is the asset; it does "
              "not leave the company. Federated learning trains one model "
              "across partners by exchanging weights, never traces. That is "
              "the promise we set out to test.")

    # 3 ─ The setting
    chart_slide(prs, 'Real partners hold real geography',
                'gap.png',
                "The setup that matters. Standard FL benchmarks shuffle data "
                "randomly across clients. Real partners each hold a contiguous "
                "region, so their facies distributions differ structurally.\n\n"
                "Read the chart left to right: centralised is the ceiling at "
                "0.693. Shuffle the data randomly and federated learning "
                "essentially matches it — 0.686. Partition it geographically, "
                "the way the world actually looks, and it falls to 0.551.\n\n"
                "Parihaka and F3, six facies, UNet, up to 20 clients.",
                takeaway='Shuffled data: federated learning works. Real geography: it does not.')

    # 4 ─ The sharpest version of the failure
    statement(prs, 'One facies went to zero.',
              'Not degraded — absent.',
              "This is the paper's sharpest result and worth a pause. Class 5, "
              "the rarest facies, reaches exactly 0.0 IoU once you have five "
              "or more clients. Not merely poor: the model stops predicting "
              "the class at all.\n\n"
              "Also worth stating: FedProx and FedBN, the standard fixes for "
              "client drift, made things worse — 1 to 12 percent. Plain FedAvg "
              "won. That tells you drift is not the problem.", accent=WARM)

    # 5 ─ Diagnosis
    statement(prs, 'The cause is absence,', 'not drift.',
              "The paper's conclusion. Most clients hold no pixels of the rare "
              "facies. Averaging their weights dilutes the few clients that do "
              "— so the class is not learned badly, it is erased. That is why "
              "drift-correction methods did not help.\n\n"
              "This is the hinge of the talk. Everything after this slide is "
              "unpublished work asking: can that erasure be undone?")

    # 6 ─ Section marker
    s = prs.slides.add_slide(prs.slide_layouts[SECTION])
    s.shapes.title.text = 'What we found since'
    s.placeholders[1].text = 'Unpublished — submitted five months ago'
    notes(s, "Signpost clearly and slow down here. Everything that follows "
             "postdates the submission.")

    # 7 ─ Bistability
    chart_slide(prs, 'The failure is bimodal, not gradual',
                'bistability.png',
                "Our most interesting finding, and it reframes the problem.\n\n"
                "Every dot is one training run — 45 of them. The rare facies "
                "either gets learned, around 0.2 to 0.3 IoU, or it collapses "
                "to exactly zero. Only three runs land anywhere in between.\n\n"
                "This is not ordinary variance; it is two attractors. Same "
                "code, same data, different random seed, opposite outcome. "
                "Method choice shifts the PROBABILITY of the good basin — it "
                "never removes the bad one.\n\n"
                "Implication: reporting a single run's number is close to "
                "meaningless here.",
                takeaway='Same code, same data, different seed — opposite outcome.')

    # 8 ─ Recovery
    chart_slide(prs, 'Recovering the gap, step by step',
                'ladder.png',
                "Four levers, each isolated, cumulative left to right.\n\n"
                "Aggregation: weight clients by how much rare-facies data they "
                "hold and how well they have learned it. Ensemble: average "
                "several seeds — which works precisely BECAUSE of the "
                "bistability we just saw; you are buying lottery tickets. "
                "Logit adjustment: a test-time correction, next slide.\n\n"
                "Together they close roughly 39 percent of the gap to "
                "centralised. Be honest that the ensemble step costs five to "
                "ten times the training compute.",
                takeaway='About 39% of the gap to centralised, recovered.')

    # 9 ─ The mechanism
    statement(prs, 'Subtract the class prior.',
              'Three lines. No retraining.',
              "Logit adjustment — MetaFusion, Chan 2019; Menon 2021. At "
              "inference, before the argmax, subtract tau times the log class "
              "prior from the logits. Frequent classes are penalised, so rare "
              "classes win more pixels.\n\n"
              "Three lines of code. No retraining. Works on any checkpoint you "
              "already have.")

    # 10 ─ The money slide
    chart_slide(prs, 'A collapsed facies, brought back',
                'rescue.png',
                "THE money slide. Say the numbers slowly, then stop talking "
                "for a beat.\n\n"
                "This is a checkpoint that had completely lost the rare "
                "facies — 0.008 IoU, effectively zero. One line of inference-"
                "time arithmetic brings it to 0.060. Seven and a half times. "
                "And mean IoU improves slightly at the same time, so it is not "
                "a trade.\n\n"
                "The headline failure from our paper is partly reversible, at "
                "zero training cost.",
                takeaway='The published failure is partly reversible — for free.')

    # 11 ─ Nuance
    chart_slide(prs, 'How hard to push depends on the model',
                'tau.png',
                "Important nuance, and it is what makes this a finding rather "
                "than a trick.\n\n"
                "Tau is not a constant. On a checkpoint that is merely "
                "majority-biased, a gentle correction — tau 0.5 — is best. On "
                "a checkpoint where the class has collapsed, you need tau 1.5. "
                "Push a healthy model that hard and you damage it.\n\n"
                "Practically: pick tau on a single held-out crossline. "
                "Essentially free.",
                takeaway='Gentle for a healthy model. Aggressive for a collapsed one.')

    # 12 ─ Self-correction
    chart_slide(prs, 'Where we had been fooling ourselves',
                'selection_bias.png',
                "Deliberate credibility slide. Volunteering this is far "
                "stronger than being asked it in Q&A.\n\n"
                "We were saving the best round by test mIoU and reporting that "
                "number. That is optimistic — and, worse, unequally so: plus "
                "0.018 for the baseline, plus 0.043 for our own best "
                "configuration.\n\n"
                "Corrected to final round, two of our claims flip. Our 'most "
                "stable' configuration is actually the least stable — one seed "
                "ends at 0.39 — and a simpler variant beats it.\n\n"
                "Say plainly: this is what the bistability does to anyone who "
                "reports a single best number, us included.",
                takeaway='Our “most stable” configuration was the least stable.')

    # 13 ─ Negative result
    chart_slide(prs, 'It is not the normalisation',
                'norm.png',
                "A negative result, and a useful one — it answers the obvious "
                "question about our own paper.\n\n"
                "The hypothesis was reasonable: BatchNorm keeps running "
                "statistics fitted to each client's own data, and we average "
                "them across clients. That is a documented failure mode under "
                "non-IID conditions, and it would neatly explain why FedBN did "
                "nothing for us.\n\n"
                "So we tested it directly — GroupNorm, which has no "
                "cross-client statistics, everything else identical. It is a "
                "tenth of an IoU WORSE, every seed, no overlap between the two "
                "sets. And it does not rescue the rare class either.\n\n"
                "Why: our clients all hold slices of one survey, so the "
                "low-level statistics — amplitude, texture, frequency — are "
                "shared. BatchNorm exploits that. The non-IID-ness here is in "
                "WHICH FACIES APPEAR, not in the input distribution. That is "
                "the distinction the FL literature does not make.\n\n"
                "So: normalisation is not the lever. The absence of data is "
                "still the problem.",
                takeaway='We tested the obvious fix directly. It is not the answer.')

    # 14 ─ Takeaways
    s = prs.slides.add_slide(prs.slide_layouts[TITLE_CONTENT])
    s.shapes.title.text = 'Three things to take away'
    body = s.placeholders[1]
    # Keep the bullet block clear of the logo strip in the bottom-right.
    body.height = Inches(4.4)
    body.width = prs.slide_width - Inches(4.4)
    tf = body.text_frame
    tf.clear()
    for i, t in enumerate(['Geography breaks federated seismic learning',
                           'Rare-class failure is bistable, not noisy',
                           'Part of the damage is reversible for free']):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.runs[0].font.size = Pt(30)
    notes(s, "Tell them what you told them.\n\n"
             "One: the realistic setting, not the benchmark setting, is where "
             "federated learning breaks.\n"
             "Two: rare-class failure is an attractor problem, not a variance "
             "problem — which changes how anyone should report results here.\n"
             "Three: some of the damage undoes at inference, for free.\n\n"
             "Then invite questions.")

    # 15 ─ Q&A
    s = prs.slides.add_slide(prs.slide_layouts[SECTION])
    s.shapes.title.text = 'Thank you'
    s.placeholders[1].text = 'Questions?'
    notes(s, "Anticipated questions:\n\n"
             "- 'Which mIoU?' Macro, unweighted, six classes.\n"
             "- 'Why did FedBN fail?' We tested it directly: GroupNorm is 0.10 "
             "IoU worse. Normalisation is not the lever — see that slide.\n"
             "- 'Did the rare-class aggregation help?' Honestly, no. Against a "
             "proper equal-weight control at n=5 it is within noise "
             "(0.5695 vs 0.5688).\n"
             "- 'What does ensembling cost?' 5-10x training. Logit adjustment "
             "is free by comparison.\n"
             "- 'Does this hold on F3?' New results are Parihaka only so far; "
             "F3 is next.\n"
             "- 'How do you pick tau without test data?' A held-out crossline "
             "from training.\n"
             "- 'Is the bistability a bug?' We looked. Same code and data, "
             "seed alone flips it — it is the optimisation landscape.")

    # Branding last, so it sits above the content on every slide. The title
    # slide gets the large treatment; content slides a discreet corner mark.
    for i, slide in enumerate(prs.slides):
        add_logos(prs, slide, large=(i == 0))

    prs.save(OUT)
    if os.path.exists(WORK):
        os.remove(WORK)
    return OUT


if __name__ == '__main__':
    print('wrote', build())
